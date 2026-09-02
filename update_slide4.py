import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PPT_PATH = r"resources/Review-1 TEMPLATE_MINI PROJECT CSE7102.pptx"
PPT_OUT = r"resources/PERIMETER_Review_1_Presentation.pptx"

def update_slide_4():
    prs = pptx.Presentation(PPT_PATH)
    s4 = prs.slides[3] # Slide 4
    
    # Clear text in shape 1 (the content box)
    for shape in s4.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Problem Statement Number" in text or "Category" in text or "Problem Description" in text:
                tf = shape.text_frame
                tf.clear()
                tf.word_wrap = True
                
                # Title / Category header
                p0 = tf.paragraphs[0]
                p0.text = "Category: Both (Hardware Sensor & Software)  |  Difficulty: Advanced"
                p0.font.size = Pt(13)
                p0.font.bold = True
                p0.font.color.rgb = RGBColor(15, 110, 110)
                p0.space_after = Pt(10)
                
                bullets = [
                    ("Context & The Crisis:", "In acute distress situations, victims face a cognitive panic freeze where manual phone unlocking or navigating apps is physically impossible within a 2-4 second threat window."),
                    ("Core Flaws of Existing Apps:", "Current solutions rely on passive SMS alerts to distant contacts (15-30km away), lack debounced motion filtering (causing high false alarms), and fail to alert nearby civilian first responders."),
                    ("The Proposed Solution (PERIMETER):", "An automated 5-tier safety network integrating hands-free accelerometer-debounced SOS (zero screen interaction), sub-second PostGIS 5km geofenced dispatch to nearby verified volunteers and police units, and an immutable universal case timeline.")
                ]
                
                for heading, body in bullets:
                    p = tf.add_paragraph()
                    run_h = p.add_run()
                    run_h.text = f"• {heading} "
                    run_h.font.bold = True
                    run_h.font.size = Pt(12)
                    run_h.font.color.rgb = RGBColor(17, 24, 39)
                    
                    run_b = p.add_run()
                    run_b.text = body
                    run_b.font.bold = False
                    run_b.font.size = Pt(11.5)
                    run_b.font.color.rgb = RGBColor(51, 65, 85)
                    p.space_after = Pt(8)

    prs.save(PPT_PATH)
    prs.save(PPT_OUT)
    print("Slide 4 updated with crisp, slide-fitting content!")

if __name__ == "__main__":
    update_slide_4()
